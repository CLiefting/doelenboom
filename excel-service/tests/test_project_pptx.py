"""Tests voor app/project_pptx.py -- bouwt de Project-PowerPoint-rapportage
voor één project en controleert de slide-structuur/inhoud door het resultaat
weer in te laden met python-pptx zelf."""
from __future__ import annotations

import io

from pptx import Presentation

from app.project_pptx import _fmt_date, build_project_pptx

from .test_project_workbook import make_data, make_meta


def all_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return '\n'.join(parts)


class TestFmtDate:
    def test_geldige_datum(self):
        assert _fmt_date('2026-09-15') == '15 sep 2026'

    def test_leeg_of_none(self):
        assert _fmt_date(None) == '-'
        assert _fmt_date('') == '-'

    def test_onherkenbare_waarde_blijft_zichtbaar(self):
        assert _fmt_date('onzin') == 'onzin'


class TestBuildProjectPptx:
    def test_bouwt_vier_slides(self):
        content = build_project_pptx(make_data(), make_meta())
        prs = Presentation(io.BytesIO(content))
        assert len(prs.slides) == 4

    def test_slide_1_toont_projectnaam_code_en_status(self):
        content = build_project_pptx(make_data(), make_meta())
        prs = Presentation(io.BytesIO(content))
        text = all_text(prs.slides[0])
        assert 'Sweepen' in text
        assert 'NP37' in text
        assert 'Groen' in text  # RAG-label, title-cased vanaf 'groen'
        assert 'actief' in text.lower()  # fixture levert lowercase 'actief'

    def test_slide_2_toont_deliverable_tabel_en_business_value(self):
        content = build_project_pptx(make_data(), make_meta())
        prs = Presentation(io.BytesIO(content))
        text = all_text(prs.slides[1])
        # Adviesrapport is nog niet opgeleverd (geen werkelijkeDatum) en heeft
        # dus een rij in de "eerstvolgende deliverables"-tabel; PID is al
        # opgeleverd (werkelijkeDatum gezet) en hoort daar niet meer in thuis.
        assert 'Adviesrapport' in text
        assert 'PID' not in text
        assert '1 van 3 opgeleverd' in text

    def test_slide_2_toont_projecttijdlijn_met_markers_en_legenda(self):
        # make_data()'s producten hebben allemaal een verwachte/werkelijke
        # datum of deadline, dus hoort de tijdlijn (as + 'vandaag'-lijn +
        # legenda) getekend te worden -- geteld via de vorm van de shapes
        # (cirkel/ruit/driehoek), niet via tekst (die staat alleen in de
        # legenda-labels en maandkoppen, niet los per marker).
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        content = build_project_pptx(make_data(), make_meta())
        prs = Presentation(io.BytesIO(content))
        slide = prs.slides[1]
        text = all_text(slide)
        assert 'vandaag' in text
        assert 'Deliverable · verwacht' in text
        assert 'Deliverable · opgeleverd' in text
        assert 'Mijlpaal · gehaald' in text
        assert 'Deadline' in text
        assert 'Sep 2026' in text or 'Okt 2026' in text  # maandkoppen onder de as

        auto_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        assert len(auto_shapes) > 5  # as-lijn, vandaag-lijn, markers, legenda-iconen

    def test_geen_gedateerde_producten_geeft_geen_tijdlijn(self):
        data = make_data()
        for p in data['products']:
            p['verwachteDatum'] = None
            p['werkelijkeDatum'] = None
            p['deadline'] = None
        content = build_project_pptx(data, make_meta())
        prs = Presentation(io.BytesIO(content))
        text = all_text(prs.slides[1])
        assert 'vandaag' not in text
        # De tabel moet dan gewoon op zijn oorspronkelijke, hogere positie
        # blijven staan -- geen lege ruimte waar de tijdlijn had gestaan.
        assert 'GO/NO-GO' in text

    def test_slide_3_toont_activiteiten_in_de_juiste_categorie(self):
        # make_data()'s activiteiten (2026-08-01 t/m 2026-08-11) liggen beide
        # vóór meta.exportedAt (2026-08-27) -- horen dus in "recent afgerond"
        # (binnen het venster van 30 dagen), niet in "loopt nu"/"gepland".
        content = build_project_pptx(make_data(), make_meta())
        prs = Presentation(io.BytesIO(content))
        text = all_text(prs.slides[2])
        assert 'Taak A' in text
        assert 'Taak B' in text
        assert 'RECENT AFGEROND' in text

    def test_slide_4_toont_toelichting_tags_en_organisatieonderdelen(self):
        content = build_project_pptx(make_data(), make_meta())
        prs = Presentation(io.BytesIO(content))
        text = all_text(prs.slides[3])
        assert 'Op schema' in text
        assert 'IGO' in text
        assert 'HRB-S' in text

    def test_lege_data_crasht_niet(self):
        content = build_project_pptx({}, {})
        prs = Presentation(io.BytesIO(content))
        assert len(prs.slides) == 4

    def test_veel_deliverables_wordt_afgekapt_met_teller(self):
        data = make_data()
        data['products'] = [
            {
                'id': i, 'name': f'D{i}', 'type': 'deliverable', 'pctGereed': 0,
                'verwachteDatum': f'2026-09-{(i % 27) + 1:02d}', 'werkelijkeDatum': None,
                'businessValue': None, 'omschrijving': '', 'deadline': None, 'duur': None,
                'duurEenheid': 'd', 'opmerking': '',
            }
            for i in range(1, 11)
        ]
        content = build_project_pptx(data, make_meta())
        prs = Presentation(io.BytesIO(content))
        text = all_text(prs.slides[1])
        assert '+ 4 andere' in text
