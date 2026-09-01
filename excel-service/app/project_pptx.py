"""
Export van de gegevens van ÉÉN project (Project-element) als PowerPoint-
presentatie -- vier slides (status, voortgang/deliverables, activiteiten,
aandachtspunten), bedoeld als kant-en-klare rapportage voor een klant/externe
stakeholder buiten de applicatie. Hergebruikt exact dezelfde 'data'/'meta'-
vorm als build_project_workbook() in project_workbook.py hiernaast (zie de
toelichting daar, en api/src/routes/projectExcel.ts::buildProjectExportData
voor waar die JSON vandaan komt) -- alleen de output is hier een .pptx
i.p.v. een .xlsx.

Puur een export, geen import/round-trip: dit is een leesbaar eindresultaat,
geen brondocument om later weer in te lezen (in tegenstelling tot het
Excel-formaat hiernaast).

Aangeroepen door api/src/routes/projectExcel.ts (POST .../project-pptx).
"""
from __future__ import annotations

import io
from datetime import date, timedelta
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Zelfde RAG-kleuren als RAG_COLORS in web/public/tree.html, voor visuele
# consistentie tussen de app zelf en dit gegenereerde document.
RAG_COLORS = {
    'rood': RGBColor(0xDC, 0x35, 0x45),
    'oranje': RGBColor(0xFD, 0x7E, 0x14),
    'groen': RGBColor(0x28, 0xA7, 0x45),
}
RAG_DEFAULT_COLOR = RGBColor(0xB5, 0xBA, 0xC2)

DARK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6C, 0x6F, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF5, 0xF7)
ACCENT = RGBColor(0x2F, 0x55, 0x97)  # zelfde blauw als de "Strategisch doel"-kolom elders in de app

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN

MAANDEN_KORT = ['jan', 'feb', 'mrt', 'apr', 'mei', 'jun', 'jul', 'aug', 'sep', 'okt', 'nov', 'dec']

MAX_ITEMS_PER_LIJST = 6


def _fmt_date(value: Any) -> str:
    """'2026-09-15' -> '15 sep 2026'; None/leeg/onherkenbaar -> '-'."""
    if not value:
        return '-'
    s = str(value)[:10]
    try:
        y, m, d = (int(part) for part in s.split('-'))
        return f'{d} {MAANDEN_KORT[m - 1]} {y}'
    except (ValueError, IndexError):
        return s


def _today_iso(meta: dict[str, Any]) -> str:
    # 'Vandaag' voor het indelen van activiteiten in "loopt nu"/"gepland" —
    # meta.exportedAt (het moment van genereren) is hier leidend i.p.v. de
    # servertijd zelf, zodat een handmatig later gedraaide her-export met een
    # meegegeven exportedAt reproduceerbaar blijft.
    exported_at = meta.get('exportedAt')
    if isinstance(exported_at, str) and exported_at:
        return exported_at[:10]
    return date.today().isoformat()


def _add_days_iso(iso_date: str, days: int) -> str:
    try:
        y, m, d = (int(part) for part in iso_date[:10].split('-'))
        return (date(y, m, d) + timedelta(days=days)).isoformat()
    except (ValueError, IndexError):
        return iso_date


def _pct(value: Any) -> int:
    try:
        return max(0, min(100, int(value)))
    except (TypeError, ValueError):
        return 0


def _num(value: Any) -> float | None:
    if value is None or value == '':
        return None
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = volledig leeg layout


def _add_rect(slide, left, top, width, height, fill: RGBColor | None, line: bool = False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = RGBColor(0xE0, 0xE2, 0xE6)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _add_text(
    slide, left, top, width, height, text: str, *, size: int, bold: bool = False,
    color: RGBColor = DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap: bool = True,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return box


def _add_bullets(
    slide, left, top, width, height, lines: list[str], *, size: int = 14,
    color: RGBColor = DARK, bullet: str = '•  ',
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = bullet + line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = 'Calibri'
    return box


def _slide_header(slide, kicker: str, title: str):
    _add_text(slide, MARGIN, Inches(0.45), CONTENT_W, Inches(0.3), kicker.upper(), size=12, bold=True, color=ACCENT)
    _add_text(slide, MARGIN, Inches(0.75), CONTENT_W, Inches(0.6), title, size=26, bold=True, color=DARK)
    _add_rect(slide, MARGIN, Inches(1.35), CONTENT_W, Emu(1), MUTED)


def _footer(slide, project: dict[str, Any], meta: dict[str, Any], page: int):
    text = (
        f"{project.get('name', '')} ({project.get('code', '')}) — "
        f"{meta.get('doelenboom', '')} / {meta.get('tenant', '')}"
    )
    _add_text(slide, MARGIN, SLIDE_H - Inches(0.4), CONTENT_W - Inches(0.6), Inches(0.3), text, size=9, color=MUTED)
    _add_text(
        slide, SLIDE_W - MARGIN - Inches(0.6), SLIDE_H - Inches(0.4), Inches(0.6), Inches(0.3),
        str(page), size=9, color=MUTED, align=PP_ALIGN.RIGHT,
    )


def _rag_color(rag: str | None) -> RGBColor:
    return RAG_COLORS.get((rag or '').strip().lower(), RAG_DEFAULT_COLOR)


# ---- Projecttijdlijn (verwachte/werkelijke opleverdatum + deadline per
# product, op één gezamenlijke maand-/kwartaalas) -- hetzelfde concept als
# productTimelineHtml/buildTimelineMarkers/timelineBandBoundaries in
# web/public/tree.html, hier eenmalig gerenderd als vaste tekening i.p.v.
# interactieve HTML (geen hover-tooltips dus, alleen de as/markers/legenda).

TIMELINE_MARKER_COLOR = RGBColor(0x2F, 0x55, 0x97)  # zelfde blauw als timelineLegendIcon in tree.html
TIMELINE_DEADLINE_COLOR = RGBColor(0xB4, 0x23, 0x18)  # zelfde rood als timelineDeadlineIcon in tree.html
TIMELINE_AXIS_COLOR = RGBColor(0xC7, 0xCB, 0xD1)


def _parse_iso_date(value: Any) -> date | None:
    if not value:
        return None
    s = str(value)[:10]
    try:
        y, m, d = (int(part) for part in s.split('-'))
        return date(y, m, d)
    except (ValueError, IndexError):
        return None


def _build_timeline_markers(products: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Zelfde opzet als buildTimelineMarkers() in tree.html: per product een
    marker voor de verwachte datum, de werkelijke (opgeleverde) datum en de
    deadline -- elk optioneel, een product kan dus 0 tot 3 markers leveren."""
    markers: list[dict[str, Any]] = []
    for p in products:
        marker_type = 'mijlpaal' if p.get('type') == 'mijlpaal' else 'deliverable'
        verwacht = _parse_iso_date(p.get('verwachteDatum'))
        if verwacht:
            markers.append({'t': verwacht, 'type': marker_type, 'filled': False, 'is_deadline': False})
        werkelijk = _parse_iso_date(p.get('werkelijkeDatum'))
        if werkelijk:
            markers.append({'t': werkelijk, 'type': marker_type, 'filled': True, 'is_deadline': False})
        deadline = _parse_iso_date(p.get('deadline'))
        if deadline:
            markers.append({'t': deadline, 'type': marker_type, 'filled': False, 'is_deadline': True})
    markers.sort(key=lambda m: m['t'])
    return markers


def _add_months(d: date, months: int) -> date:
    total = d.year * 12 + (d.month - 1) + months
    y, m = divmod(total, 12)
    return date(y, m + 1, 1)


def _timeline_bounds(markers: list[dict[str, Any]], today: date) -> tuple[bool, list[date]]:
    """Zelfde as-logica als timelineBandBoundaries()/computeProjectTimelineBounds()
    in tree.html: 'vandaag' telt altijd mee in het bereik, bij een spanne van
    meer dan ~460 dagen worden het kwartalen i.p.v. maanden."""
    all_dates = [m['t'] for m in markers] + [today]
    raw_min, raw_max = min(all_dates), max(all_dates)
    if raw_min == raw_max:
        raw_min -= timedelta(days=1)
        raw_max += timedelta(days=1)
    quarterly = (raw_max - raw_min).days > 460
    step = 3 if quarterly else 1
    first = date(raw_min.year, raw_min.month, 1)
    if quarterly:
        first = date(first.year, ((first.month - 1) // 3) * 3 + 1, 1)
    bounds = [first]
    while bounds[-1] < raw_max:
        bounds.append(_add_months(bounds[-1], step))
    return quarterly, bounds


def _add_project_timeline(slide, top, products: list[dict[str, Any]], today_iso: str):
    """Tekent de projecttijdlijn en geeft de Y-positie net onder de tijdlijn
    terug, zodat de aanroeper de inhoud eronder kan plaatsen -- of None als
    geen enkel product een verwachte/werkelijke datum of deadline heeft (dan
    is er niets te plotten, zelfde als productTimelineHtml() '' in tree.html)."""
    markers = _build_timeline_markers(products)
    if not markers:
        return None

    today = _parse_iso_date(today_iso) or date.today()
    quarterly, bounds = _timeline_bounds(markers, today)
    axis_start, axis_end = bounds[0], bounds[-1]
    span_days = (axis_end - axis_start).days or 1

    # Kleine marge aan weerszijden zodat een marker precies op het begin/eind
    # van het bereik (bv. een mijlpaal exact op de laatste maandgrens) niet
    # half buiten de tijdlijn/slide valt.
    pad = Inches(0.1)
    inner_left = MARGIN + pad
    inner_w = CONTENT_W - 2 * pad

    def x_for(d: date) -> int:
        frac = (d - axis_start).days / span_days
        return int(inner_left + inner_w * frac)

    axis_y = top + Inches(0.68)

    for i in range(len(bounds) - 1):
        left = x_for(bounds[i])
        width = x_for(bounds[i + 1]) - left
        if width >= Inches(0.75):
            label = (
                f'K{(bounds[i].month - 1) // 3 + 1} {bounds[i].year}' if quarterly
                else f'{MAANDEN_KORT[bounds[i].month - 1].capitalize()} {bounds[i].year}'
            )
            _add_text(slide, left, axis_y + Inches(0.08), width, Inches(0.25), label, size=9, color=MUTED, align=PP_ALIGN.CENTER)

    _add_rect(slide, MARGIN, axis_y, CONTENT_W, Pt(1.25), TIMELINE_AXIS_COLOR)

    if axis_start <= today <= axis_end:
        today_x = x_for(today)
        _add_rect(slide, today_x, top, Pt(1.25), Inches(0.62), ACCENT)
        _add_text(slide, today_x - Inches(0.35), top - Inches(0.02), Inches(0.7), Inches(0.2), 'vandaag', size=8, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Eenvoudige verticale stapeling om markers die (bijna) op dezelfde datum
    # vallen niet exact over elkaar te laten landen -- zelfde bucket-aanpak
    # als productTimelineHtml in tree.html (geen echte collision-detectie,
    # maar volstaat voor de doorgaans kleine aantallen items per project).
    marker_size = Inches(0.14)
    half = Inches(0.07)
    bucket_counts: dict[int, int] = {}
    for m in markers:
        cx = x_for(m['t'])
        bucket = round((cx - MARGIN) / CONTENT_W * 60)
        stack = bucket_counts.get(bucket, 0)
        bucket_counts[bucket] = stack + 1
        level = stack % 3
        cy = axis_y - Inches(0.14) - int(Inches(0.16) * level)
        if m['is_deadline']:
            shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, cx - half, cy - half, marker_size, marker_size)
            shape.rotation = 180
            shape.fill.solid()
            shape.fill.fore_color.rgb = TIMELINE_DEADLINE_COLOR
            shape.line.color.rgb = TIMELINE_DEADLINE_COLOR
        else:
            mso = MSO_SHAPE.DIAMOND if m['type'] == 'mijlpaal' else MSO_SHAPE.OVAL
            shape = slide.shapes.add_shape(mso, cx - half, cy - half, marker_size, marker_size)
            shape.line.color.rgb = TIMELINE_MARKER_COLOR
            shape.line.width = Pt(1.25)
            shape.fill.solid()
            shape.fill.fore_color.rgb = TIMELINE_MARKER_COLOR if m['filled'] else WHITE
        shape.shadow.inherit = False

    legend_y = axis_y + Inches(0.38)
    legend_items = [
        ('deliverable', False, False, 'Deliverable · verwacht'),
        ('deliverable', True, False, 'Deliverable · opgeleverd'),
        ('mijlpaal', False, False, 'Mijlpaal · verwacht'),
        ('mijlpaal', True, False, 'Mijlpaal · gehaald'),
        (None, False, True, 'Deadline'),
    ]
    dot_size = Inches(0.11)
    legend_col_w = Inches(2.1)
    legend_x = MARGIN
    for marker_type, filled, is_deadline, label in legend_items:
        if is_deadline:
            shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, legend_x, legend_y, dot_size, dot_size)
            shape.rotation = 180
            shape.fill.solid()
            shape.fill.fore_color.rgb = TIMELINE_DEADLINE_COLOR
            shape.line.color.rgb = TIMELINE_DEADLINE_COLOR
        else:
            mso = MSO_SHAPE.DIAMOND if marker_type == 'mijlpaal' else MSO_SHAPE.OVAL
            shape = slide.shapes.add_shape(mso, legend_x, legend_y, dot_size, dot_size)
            shape.line.color.rgb = TIMELINE_MARKER_COLOR
            shape.line.width = Pt(1)
            shape.fill.solid()
            shape.fill.fore_color.rgb = TIMELINE_MARKER_COLOR if filled else WHITE
        shape.shadow.inherit = False
        _add_text(slide, legend_x + dot_size + Inches(0.08), legend_y - Inches(0.02), legend_col_w - dot_size - Inches(0.08), Inches(0.25), label, size=9, color=MUTED)
        legend_x += legend_col_w

    return legend_y + Inches(0.35)


def _slide_status(prs: Presentation, project: dict[str, Any], meta: dict[str, Any]):
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_BG)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.18), ACCENT)

    _add_text(
        slide, MARGIN, Inches(0.9), CONTENT_W, Inches(0.35),
        (meta.get('doelenboom') or '').upper() + ('  ·  ' + (project.get('code') or '') if project.get('code') else ''),
        size=13, bold=True, color=ACCENT,
    )
    _add_text(slide, MARGIN, Inches(1.3), CONTENT_W, Inches(1.2), project.get('name') or 'Project', size=40, bold=True, color=DARK)

    description = (project.get('description') or '').strip()
    if description:
        _add_text(slide, MARGIN, Inches(2.35), CONTENT_W, Inches(0.9), description, size=15, color=MUTED)

    status = project.get('status') or {}
    rag = status.get('rag') or ''
    projectstatus = status.get('projectstatus') or 'Onbekend'
    rag_label = rag.title() if rag else 'Niet gerapporteerd'

    badge_w, badge_h = Inches(3.4), Inches(1.3)
    badge_top = Inches(3.6)
    _add_rect(slide, MARGIN, badge_top, badge_w, badge_h, _rag_color(rag))
    _add_text(
        slide, MARGIN, badge_top + Inches(0.18), badge_w, Inches(0.5), rag_label,
        size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide, MARGIN, badge_top + Inches(0.72), badge_w, Inches(0.4), 'RAG-status',
        size=12, color=WHITE, align=PP_ALIGN.CENTER,
    )

    info_left = MARGIN + badge_w + Inches(0.5)
    info_w = CONTENT_W - badge_w - Inches(0.5)
    _add_text(slide, info_left, badge_top, info_w, Inches(0.35), 'PROJECTSTATUS', size=11, bold=True, color=MUTED)
    _add_text(slide, info_left, badge_top + Inches(0.32), info_w, Inches(0.5), projectstatus, size=22, bold=True, color=DARK)
    _add_text(
        slide, info_left, badge_top + Inches(0.9), info_w, Inches(0.35),
        'Gerapporteerd op ' + _fmt_date(status.get('gerapporteerdOp')), size=13, color=MUTED,
    )

    _footer(slide, project, meta, 1)


def _slide_voortgang(prs: Presentation, project: dict[str, Any], products: list[dict[str, Any]], meta: dict[str, Any]):
    slide = _blank_slide(prs)
    _slide_header(slide, 'Voortgang', 'Oplevering & deliverables')

    total = len(products)
    delivered = sum(1 for p in products if p.get('werkelijkeDatum'))
    weighted_bv = 0.0
    total_bv = 0.0
    for p in products:
        bv = _num(p.get('businessValue'))
        if bv is None:
            continue
        total_bv += bv
        weighted_bv += bv * (_pct(p.get('pctGereed')) / 100)

    summary_bits = [f'{delivered} van {total} opgeleverd']
    if total_bv > 0:
        summary_bits.append(f'business value {round(weighted_bv)} / {round(total_bv)} gerealiseerd')
    _add_text(slide, MARGIN, Inches(1.55), CONTENT_W, Inches(0.4), '  •  '.join(summary_bits), size=15, bold=True, color=ACCENT)

    # Projecttijdlijn (verwachte/werkelijke opleverdatum + deadline per
    # product) -- zelfde as/markers als op de projectkaart in de app. Geeft
    # None terug als geen enkel product een datum heeft; dan blijft de tabel
    # hieronder op zijn oorspronkelijke, hogere positie staan.
    timeline_bottom = _add_project_timeline(slide, Inches(1.95), products, _today_iso(meta))

    # Eerstvolgende, nog niet opgeleverde deliverables/mijlpalen, op
    # verwachte datum -- dat is voor een externe lezer relevanter dan een
    # volledige, mogelijk lange lijst van alles wat al klaar is.
    upcoming = [p for p in products if not p.get('werkelijkeDatum')]
    upcoming.sort(key=lambda p: p.get('verwachteDatum') or '9999-99-99')
    shown = upcoming[:MAX_ITEMS_PER_LIJST]

    rows = len(shown) + 1
    table_top = timeline_bottom + Inches(0.2) if timeline_bottom else Inches(2.15)
    row_height = Inches(0.4)
    table_height = row_height * rows
    gfx = slide.shapes.add_table(rows, 4, MARGIN, table_top, CONTENT_W, table_height)
    table = gfx.table
    for row in table.rows:
        row.height = row_height
    table.columns[0].width = Inches(6.3)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.4)
    table.columns[3].width = int(CONTENT_W) - Inches(6.3) - Inches(2.2) - Inches(2.4)

    headers = ['Deliverable', 'Type', 'Verwachte datum', '% gereed']
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE

    for r, p in enumerate(shown, start=1):
        type_label = 'Mijlpaal' if p.get('type') == 'mijlpaal' else 'Deliverable'
        values = [p.get('name') or '', type_label, _fmt_date(p.get('verwachteDatum')), f"{_pct(p.get('pctGereed'))}%"]
        for c, v in enumerate(values):
            cell = table.cell(r, c)
            cell.text = v
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT_BG

    if not shown:
        _add_text(
            slide, MARGIN, table_top, CONTENT_W, Inches(0.4),
            'Geen openstaande deliverables.' if total else 'Nog geen deliverables vastgelegd voor dit project.',
            size=13, color=MUTED,
        )
    elif len(upcoming) > len(shown):
        _add_text(
            slide, MARGIN, table_top + table_height + Inches(0.1), CONTENT_W, Inches(0.3),
            f'+ {len(upcoming) - len(shown)} andere nog te leveren deliverable(s)', size=11, color=MUTED,
        )

    _footer(slide, project, meta, 2)


AFGEROND_VENSTER_DAGEN = 30


def _slide_activiteiten(prs: Presentation, project: dict[str, Any], activities: list[dict[str, Any]], meta: dict[str, Any]):
    slide = _blank_slide(prs)
    _slide_header(slide, 'Planning', 'Wat gebeurt er nu en wat komt eraan')

    today = _today_iso(meta)
    # Fase/samenvattende rijen zijn in de app inklapbare groep-headers, geen
    # losse werkitems -- die laten we hier buiten beschouwing (zelfde als
    # hoe project_workbook.py ze wél opneemt voor round-trip, maar dit is
    # geen brondocument).
    real = [a for a in activities if not a.get('isSummary')]

    def end_of(a: dict[str, Any]) -> str:
        return a.get('endDate') or a.get('startDate') or ''

    def start_of(a: dict[str, Any]) -> str:
        return a.get('startDate') or ''

    # Drie categorieën t.o.v. 'vandaag' (zie _today_iso) -- i.p.v. alleen
    # "loopt nu"/"gepland": een activiteit die net vóór het rapportagemoment
    # is afgerond hoort ook in een statusrapportage thuis, en zonder deze
    # categorie zou de slide leeg kunnen blijven als er toevallig niets exact
    # over 'vandaag' loopt.
    afgerond_grens = _add_days_iso(today, -AFGEROND_VENSTER_DAGEN)
    afgerond = [a for a in real if end_of(a) < today and end_of(a) >= afgerond_grens]
    lopend = [a for a in real if start_of(a) <= today <= end_of(a)]
    gepland = [a for a in real if start_of(a) > today]
    afgerond.sort(key=end_of, reverse=True)
    lopend.sort(key=end_of)
    gepland.sort(key=start_of)

    def line_for(a: dict[str, Any]) -> str:
        marker = '◆ ' if a.get('isMilestone') else ''
        start, end = _fmt_date(a.get('startDate')), _fmt_date(a.get('endDate'))
        when = start if a.get('isMilestone') or start == end else f'{start} — {end}'
        return f'{marker}{a.get("name") or ""}  ({when})'

    gap = Inches(0.4)
    col_w = (CONTENT_W - 2 * gap) / 3
    content_top = Inches(1.65)
    columns = [
        (f'RECENT AFGEROND ({len(afgerond)})', afgerond, 'Geen recent afgeronde activiteiten.'),
        (f'LOOPT NU ({len(lopend)})', lopend, 'Geen lopende activiteiten.'),
        (f'GEPLAND ({len(gepland)})', gepland, 'Geen geplande activiteiten.'),
    ]
    for i, (label, items, empty_text) in enumerate(columns):
        col_left = MARGIN + i * (col_w + gap)
        _add_text(slide, col_left, content_top, col_w, Inches(0.35), label, size=13, bold=True, color=ACCENT)
        if items:
            _add_bullets(
                slide, col_left, content_top + Inches(0.45), col_w, Inches(4.5),
                [line_for(a) for a in items[:MAX_ITEMS_PER_LIJST]], size=12,
            )
        else:
            _add_text(slide, col_left, content_top + Inches(0.45), col_w, Inches(0.4), empty_text, size=11, color=MUTED)

    if any(len(items) > MAX_ITEMS_PER_LIJST for _, items, _ in columns):
        _add_text(
            slide, MARGIN, SLIDE_H - Inches(0.75), CONTENT_W, Inches(0.3),
            '◆ = mijlpaal — niet alle activiteiten passen op deze slide, zie de volledige planning in de app.',
            size=10, color=MUTED,
        )

    _footer(slide, project, meta, 3)


def _slide_aandachtspunten(prs: Presentation, project: dict[str, Any], meta: dict[str, Any]):
    slide = _blank_slide(prs)
    _slide_header(slide, 'Aandachtspunten', 'Toelichting & vervolg')

    status = project.get('status') or {}
    toelichting = (status.get('toelichting') or '').strip()
    _add_text(
        slide, MARGIN, Inches(1.6), CONTENT_W, Inches(2.2),
        toelichting or 'Geen toelichting vastgelegd bij de huidige status.',
        size=16, color=DARK,
    )

    chip_top = Inches(4.1)
    tags = project.get('tags') or []
    if tags:
        _add_text(slide, MARGIN, chip_top, CONTENT_W, Inches(0.3), 'TAGS', size=11, bold=True, color=MUTED)
        _add_text(slide, MARGIN, chip_top + Inches(0.3), CONTENT_W, Inches(0.4), '  ·  '.join(tags), size=13, color=DARK)
        chip_top += Inches(0.85)

    orgs = project.get('orgs') or []
    if orgs:
        _add_text(slide, MARGIN, chip_top, CONTENT_W, Inches(0.3), 'ORGANISATIEONDERDELEN', size=11, bold=True, color=MUTED)
        org_line = '  ·  '.join(f"{o.get('name', '')} ({o.get('relatietype', '')})" for o in orgs)
        _add_text(slide, MARGIN, chip_top + Inches(0.3), CONTENT_W, Inches(0.4), org_line, size=13, color=DARK)
        chip_top += Inches(0.85)

    if status.get('clusterPpt'):
        _add_text(slide, MARGIN, chip_top, CONTENT_W, Inches(0.3), 'CLUSTER PPT', size=11, bold=True, color=MUTED)
        _add_text(slide, MARGIN, chip_top + Inches(0.3), CONTENT_W, Inches(0.4), status['clusterPpt'], size=13, color=DARK)

    generated_by = meta.get('exportedBy') or 'onbekend'
    generated_at = _fmt_date(meta.get('exportedAt'))
    _add_text(
        slide, MARGIN, SLIDE_H - Inches(0.75), CONTENT_W, Inches(0.3),
        f'Automatisch gegenereerd op {generated_at} door {generated_by} vanuit Doelenboom.',
        size=10, color=MUTED,
    )
    _footer(slide, project, meta, 4)


def build_project_pptx(data: dict[str, Any], meta: dict[str, Any]) -> bytes:
    project = data.get('project') or {}
    products = data.get('products') or []
    activities = data.get('activities') or []

    prs = _new_presentation()
    _slide_status(prs, project, meta)
    _slide_voortgang(prs, project, products, meta)
    _slide_activiteiten(prs, project, activities, meta)
    _slide_aandachtspunten(prs, project, meta)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
